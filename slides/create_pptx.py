#!/usr/bin/env python3
"""Generate CWL Workshop PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
SKA_BLUE = RGBColor(0, 163, 224)
ACCENT_ORANGE = RGBColor(255, 107, 53)
DARK_BG = RGBColor(10, 14, 23)
CARD_BG = RGBColor(20, 24, 36)
TEXT_PRIMARY = RGBColor(232, 236, 244)
TEXT_MUTED = RGBColor(122, 133, 153)

def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_accent_bar(slide):
    """Add the left accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        Inches(0.08), Inches(5.625)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SKA_BLUE
    shape.line.fill.background()

def create_presentation():
    """Create the workshop presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, DARK_BG)
    add_accent_bar(slide)
    
    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "HANDS-ON WORKSHOP"
    p.font.size = Pt(14)
    p.font.color.rgb = SKA_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Common Workflow Language"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "for SKA Data Processing"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle description
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Building Reproducible, Portable Pipelines for Radio Astronomy"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER
    
    # Info boxes
    info_y = 4.2
    for i, (label, value) in enumerate([("Duration", "3 Hours"), ("Format", "Hands-on"), ("Level", "Beginner")]):
        x = 2.5 + i * 2
        txBox = slide.shapes.add_textbox(Inches(x), Inches(info_y), Inches(1.5), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = value
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: What is CWL
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_accent_bar(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "What is CWL?"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p = tf.add_paragraph()
    p.text = "Common Workflow Language"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    
    # Content boxes
    concepts = [
        ("Open Standard", "A specification for describing analysis workflows and tools in a portable, scalable manner."),
        ("YAML-Based", "Human-readable syntax that's easy to write, version control, and share."),
        ("Container Native", "First-class Docker/Singularity support for reproducible environments.")
    ]
    
    for i, (title, desc) in enumerate(concepts):
        y = 1.4 + i * 1.2
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(5), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.15), Inches(4.6), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = SKA_BLUE
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
    
    # Write Once Run Anywhere box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(1.8), Inches(3.5), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SKA_BLUE
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(6.2), Inches(2.2), Inches(3.1), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "↻"
    p.font.size = Pt(60)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Write Once"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Run Anywhere"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 3: Why CWL for SKA
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_accent_bar(slide)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Why CWL for SKA?"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    
    stats = [
        ("710 PB", "Data per year at full operations", "Distributed processing across SRCNet requires portable, reproducible workflows."),
        ("Global", "Collaboration across continents", "Standard format enables sharing pipelines between institutions worldwide."),
        ("50+ yr", "Science archive lifetime", "Reproducibility crucial for long-term scientific integrity and reprocessing."),
        ("HPC", "Multi-platform execution", "Same workflow runs on laptop, HPC cluster, and cloud infrastructure.")
    ]
    
    for i, (num, subtitle, desc) in enumerate(stats):
        x = 0.5 + i * 2.4
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.3), Inches(2.2), Inches(3.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()
        
        # Blue left border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.3), Inches(0.04), Inches(3.8))
        border.fill.solid()
        border.fill.fore_color.rgb = SKA_BLUE
        border.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.5), Inches(1.9), Inches(3.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ORANGE
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_PRIMARY
    
    # Slide 4: Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_accent_bar(slide)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Today's Journey"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p = tf.add_paragraph()
    p.text = "3-hour hands-on workshop"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    
    agenda = [
        ("0:00", "Introduction & Setup", "Verify environments, overview of CWL", "15 min", SKA_BLUE),
        ("0:15", "CWL Fundamentals", "YAML syntax, CommandLineTool basics", "30 min", SKA_BLUE),
        ("0:45", "Exercises 1 & 2", "Hello CWL, FITS Header Extraction", "30 min", ACCENT_ORANGE),
        ("1:15", "☕ Break", "", "15 min", TEXT_MUTED),
        ("1:30", "Workflows & Exercise 3", "Multi-step pipelines, Imaging Pipeline", "45 min", ACCENT_ORANGE),
        ("2:15", "SKA Use Case & Exercise 4", "Full calibration pipeline", "30 min", ACCENT_ORANGE),
        ("2:45", "Wrap-up & Next Steps", "Resources, community, Q&A", "15 min", SKA_BLUE),
    ]
    
    for i, (time, title, desc, duration, color) in enumerate(agenda):
        y = 1.2 + i * 0.58
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.52))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()
        
        if color != TEXT_MUTED:
            border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.04), Inches(0.52))
            border.fill.solid()
            border.fill.fore_color.rgb = color
            border.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.1), Inches(0.8), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = time
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.08), Inches(6), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRIMARY if color != TEXT_MUTED else TEXT_MUTED
        if desc:
            p = tf.add_paragraph()
            p.text = desc
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_MUTED
        
        txBox = slide.shapes.add_textbox(Inches(8.5), Inches(y + 0.15), Inches(0.8), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = duration
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
    
    # Slide 5: Core Concepts
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_accent_bar(slide)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CWL Core Concepts"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p = tf.add_paragraph()
    p.text = "Two building blocks for everything"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    
    # CommandLineTool box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(4.5), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SKA_BLUE
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(4.1), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CommandLineTool"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "Wraps a single command-line program"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Workflow box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.5), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_ORANGE
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(5.4), Inches(1.4), Inches(4.1), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Workflow"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "Connects multiple tools together"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Code example boxes
    for i, (code, color) in enumerate([
        ("cwlVersion: v1.2\nclass: CommandLineTool\n\nbaseCommand: echo\n\ninputs:\n  message:\n    type: string\n    inputBinding:\n      position: 1\n\noutputs:\n  output:\n    type: stdout", SKA_BLUE),
        ("cwlVersion: v1.2\nclass: Workflow\n\ninputs:\n  input_file: File\n\nsteps:\n  step1:\n    run: tool1.cwl\n    in:\n      file: input_file\n    out: [result]\n\noutputs:\n  final:\n    outputSource: step1/result", ACCENT_ORANGE)
    ]):
        x = 0.5 + i * 4.7
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.1), Inches(4.5), Inches(3.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.2), Inches(4.1), Inches(3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "EXAMPLE"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = color
        p = tf.add_paragraph()
        p.text = code
        p.font.size = Pt(9)
        p.font.name = "Courier New"
        p.font.color.rgb = TEXT_PRIMARY
    
    # Slide 6: Exercises Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_accent_bar(slide)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Hands-on Exercises"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p = tf.add_paragraph()
    p.text = "Progressive complexity, real astronomy applications"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    
    exercises = [
        ("1", "Hello CWL", "Your first CommandLineTool. Learn YAML syntax and basic structure.", "exercises/01-hello-cwl/", SKA_BLUE),
        ("2", "FITS Header", "Docker containers and File I/O. Extract metadata from FITS files.", "exercises/02-fits-header/", SKA_BLUE),
        ("3", "Imaging Pipeline", "Multi-step workflows. Chain tools to create an analysis pipeline.", "exercises/03-imaging-pipeline/", ACCENT_ORANGE),
        ("4", "SKA Calibration", "Real-world pipeline. Flagging, calibration, imaging, and QA.", "exercises/04-ska-calibration/", ACCENT_ORANGE),
    ]
    
    for i, (num, title, desc, path, color) in enumerate(exercises):
        x = 0.5 + i * 2.4
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(2.2), Inches(3.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()
        
        # Top border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.3), Inches(2.2), Inches(0.04))
        border.fill.solid()
        border.fill.fore_color.rgb = color
        border.line.fill.background()
        
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(1.5), Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.53), Inches(0.35), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        txBox = slide.shapes.add_textbox(Inches(x + 0.6), Inches(1.5), Inches(1.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRIMARY
        
        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.0), Inches(1.9), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        
        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(4.7), Inches(1.9), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = path
        p.font.size = Pt(9)
        p.font.color.rgb = color
    
    # Slide 7: Resources
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_accent_bar(slide)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Resources & Next Steps"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p = tf.add_paragraph()
    p.text = "Your journey continues"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    
    sections = [
        ("DOCUMENTATION", [
            ("CWL Specification", "commonwl.org/specification"),
            ("CWL User Guide", "commonwl.org/user_guide"),
            ("SKA Developer Portal", "developer.skao.int"),
        ], SKA_BLUE),
        ("COMMUNITY", [
            ("CWL Discourse", "cwl.discourse.group"),
            ("Slack: #cwl-help", "Office hours: Wed 14:00 UTC"),
            ("SKA CWL Working Group", "Monthly meetings"),
        ], SKA_BLUE),
        ("ACTION ITEMS", [
            ("① Complete survey", "Help us improve!"),
            ("② Join office hours", "Next week, get 1:1 help"),
            ("③ Build your pipeline", "Start with your own data!"),
        ], ACCENT_ORANGE),
    ]
    
    for i, (header, items, color) in enumerate(sections):
        x = 0.5 + i * 3.2
        
        txBox = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(3), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color
        
        for j, (title, subtitle) in enumerate(items):
            y = 1.6 + j * 1.1
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3), Inches(0.9))
            shape.fill.solid()
            shape.fill.fore_color.rgb = CARD_BG
            shape.line.fill.background()
            
            if color == ACCENT_ORANGE:
                border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.03), Inches(0.9))
                border.fill.solid()
                border.fill.fore_color.rgb = color
                border.line.fill.background()
            
            txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(2.7), Inches(0.7))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_PRIMARY
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_MUTED
    
    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), 'cwl-workshop-slides.pptx')
    prs.save(output_path)
    print(f"Created: {output_path}")

if __name__ == "__main__":
    create_presentation()
